#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Crafter DQN project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add authors
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
    author_frame = author_box.text_frame
    author_frame.text = "[Your Names Here]"
    author_frame.paragraphs[0].font.size = Pt(20)
    author_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "November 2025"
    date_frame.paragraphs[0].font.size = Pt(18)
    date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add key results box
    results_box = slide.shapes.add_textbox(Inches(1.5), Inches(6), Inches(7), Inches(1))
    results_frame = results_box.text_frame
    results_frame.text = "✓ 4.90 reward vs -0.90 random (5.4× better)  •  ✓ 10/22 achievements  •  ✓ Double DQN + Dueling = optimal"
    results_frame.paragraphs[0].font.size = Pt(14)
    results_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    results_frame.paragraphs[0].font.italic = True

    return slide

def add_content_slide(prs, title):
    """Add content slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_bullet_text(text_frame, text, level=0, font_size=18, bold=False):
    """Add bullet point to text frame"""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(font_size)
    p.font.bold = bold
    return p

# Slide 1: Title
print("Creating Slide 1: Title...")
add_title_slide(prs,
    "Deep Q-Network Agent for Crafter Environment",
    "Systematic Investigation of DQN Enhancements and Exploration Strategies")

# Slide 2: The Challenge
print("Creating Slide 2: Challenge...")
slide = add_content_slide(prs, "The Crafter Challenge")
left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.5))
tf = left_box.text_frame
tf.word_wrap = True

add_bullet_text(tf, "Environment Specifications:", level=0, font_size=20, bold=True)
add_bullet_text(tf, "64×64×3 RGB observations", level=1, font_size=16)
add_bullet_text(tf, "17 discrete actions", level=1, font_size=16)
add_bullet_text(tf, "22 achievements", level=1, font_size=16)
add_bullet_text(tf, "1M step budget", level=1, font_size=16)
add_bullet_text(tf, "Procedurally generated maps", level=1, font_size=16)

tf.add_paragraph()
add_bullet_text(tf, "Key Challenges:", level=0, font_size=20, bold=True)
add_bullet_text(tf, "Sparse Rewards", level=1, font_size=16)
add_bullet_text(tf, "Credit Assignment", level=1, font_size=16)
add_bullet_text(tf, "Exploration", level=1, font_size=16)
add_bullet_text(tf, "Long Episodes (150-200 steps)", level=1, font_size=16)

# Right box - baseline
right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4), Inches(2))
tf2 = right_box.text_frame
tf2.word_wrap = True
add_bullet_text(tf2, "Random Baseline:", level=0, font_size=18, bold=True)
add_bullet_text(tf2, "-0.90 ± 0.00 reward", level=1, font_size=16)
add_bullet_text(tf2, "0/22 achievements", level=1, font_size=16)
add_bullet_text(tf2, "0% Crafter Score", level=1, font_size=16)

# Slide 3: Method
print("Creating Slide 3: Method...")
slide = add_content_slide(prs, "Method: Enhanced DQN Architecture")
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf = text_box.text_frame
tf.word_wrap = True

add_bullet_text(tf, "Base DQN Loss Function:", level=0, font_size=18, bold=True)
add_bullet_text(tf, "L(θ) = E[(r + γ · max Q_target(s', a') - Q(s, a))²]", level=1, font_size=14)

tf.add_paragraph()
add_bullet_text(tf, "Enhancement 1: Double DQN", level=0, font_size=18, bold=True)
add_bullet_text(tf, "L(θ) = E[(r + γ · Q_target(s', argmax Q(s', a')) - Q(s, a))²]", level=1, font_size=14)
add_bullet_text(tf, "Why? Decouples action selection from evaluation → reduces overestimation bias", level=1, font_size=14)

tf.add_paragraph()
add_bullet_text(tf, "Enhancement 2: Dueling Architecture", level=0, font_size=18, bold=True)
add_bullet_text(tf, "Q(s, a) = V(s) + [A(s, a) - mean A(s, a')]", level=1, font_size=14)
add_bullet_text(tf, "Why? Separates state value from action advantage → better value estimates", level=1, font_size=14)

tf.add_paragraph()
add_bullet_text(tf, "Enhancement 3: N-Step Returns (n=1, 3, 5)", level=0, font_size=18, bold=True)
add_bullet_text(tf, "R_t^(n) = Σ γ^i · r_{t+i} + γ^n · max Q(s_{t+n}, a)", level=1, font_size=14)
add_bullet_text(tf, "Why? Faster propagation of sparse rewards → better credit assignment", level=1, font_size=14)

tf.add_paragraph()
add_bullet_text(tf, "Hyperparameters: LR=1e-4, Buffer=100k, Batch=32, Target Update=2500 steps", level=0, font_size=12)

# Slide 4: Ablation Study
print("Creating Slide 4: Ablation Study...")
slide = add_content_slide(prs, "Ablation Study: Which Enhancements Matter?")

# Add table
rows, cols = 7, 4
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(4.5)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(3)
table.columns[1].width = Inches(2)
table.columns[2].width = Inches(2)
table.columns[3].width = Inches(2)

# Header row
headers = ["Configuration", "Final Reward", "Crafter Score", "Achievements"]
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(14)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Data rows
data = [
    ["Random Baseline", "-0.90 ± 0.00", "0.00%", "0.0 ± 0.0"],
    ["Base DQN", "4.17 ± 0.10", "30.65%", "8.3 ± 0.5"],
    ["+ Double + Dueling ⭐", "4.90 ± 0.16", "32.91%", "8.0 ± 0.0"],
    ["+ N-step (n=3)", "4.48 ± 0.16", "30.23%", "8.7 ± 0.9"],
    ["+ N-step (n=5)", "4.00 ± 0.19", "26.82%", "8.7 ± 1.2"],
    ["All Enhancements*", "2.37 ± 0.06", "43.70%", "4.7 ± 0.5"],
]

for row_idx, row_data in enumerate(data, start=1):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_text
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        if "⭐" in cell_text:
            cell.text_frame.paragraphs[0].font.bold = True

# Add note
note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(1))
tf = note_box.text_frame
tf.text = "Note: All results averaged over 3 seeds (0, 1, 2) at 1M steps. *Includes Munchausen-DQN (showed instability)"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.italic = True

# Slide 5: Exploration Experiment
print("Creating Slide 5: Exploration Experiment...")
slide = add_content_slide(prs, "Exploration Experiment: Does Longer Help?")
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(2))
tf = text_box.text_frame
tf.word_wrap = True

add_bullet_text(tf, "Hypothesis: Extended epsilon-greedy exploration will help agent discover more strategies", level=0, font_size=16, bold=True)
add_bullet_text(tf, "Standard Schedule: Decay ε from 1.0 to 0.01 over 50,000 steps (5% of training)", level=1, font_size=14)
add_bullet_text(tf, "Extended Schedule: Decay ε from 1.0 to 0.05 over 200,000 steps (20% of training)", level=1, font_size=14)

# Results table
rows, cols = 4, 4
top = Inches(3.5)
table = slide.shapes.add_table(rows, cols, Inches(0.5), top, Inches(9), Inches(2.5)).table

# Headers
headers = ["Configuration", "Standard (50k ε)", "Extended (200k ε)", "Change"]
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(12)

# Data
data = [
    ["Base DQN", "4.17 ± 0.10 ⭐", "3.53 ± 0.77", "-15% ❌"],
    ["Enhanced DQN", "4.90 ± 0.16 ⭐", "4.33 ± 0.16", "-12% ❌"],
    ["Key Finding", "Well-calibrated", "Higher variance (7.7×)", "Longer ≠ Better"],
]

for row_idx, row_data in enumerate(data, start=1):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_text
        cell.text_frame.paragraphs[0].font.size = Pt(11)

# Add finding box
finding_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.8))
tf = finding_box.text_frame
tf.text = "Finding: Extended exploration DECREASED performance. Standard 50k decay is well-calibrated!"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(204, 0, 0)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 6: Training Dynamics (with plot)
print("Creating Slide 6: Training Dynamics...")
slide = add_content_slide(prs, "Training Dynamics: Best Agent Performance")

# Add plot
plot_path = "logdir/enhanced-dqn/plots/rewards.png"
if os.path.exists(plot_path):
    left = Inches(0.5)
    top = Inches(1.5)
    pic = slide.shapes.add_picture(plot_path, left, top, width=Inches(9))
    print(f"  Added plot: {plot_path}")
else:
    print(f"  WARNING: Plot not found: {plot_path}")

# Add results box
results_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.8))
tf = results_box.text_frame
tf.word_wrap = True
add_bullet_text(tf, "Enhanced DQN: 4.90 ± 0.16 reward (avg over 3 seeds) | Random: -0.90", level=0, font_size=14, bold=True)
add_bullet_text(tf, "MANDATORY PLOT: Average episodic reward showing training & evaluation performance", level=0, font_size=11)

# Slide 7: Achievement Spectrum (with plot)
print("Creating Slide 7: Achievement Spectrum...")
slide = add_content_slide(prs, "Achievement Spectrum Analysis (BONUS)")

# Add plot
plot_path = "logdir/enhanced-dqn/plots/achievement_spectrum.png"
if os.path.exists(plot_path):
    left = Inches(1)
    top = Inches(1.5)
    pic = slide.shapes.add_picture(plot_path, left, top, width=Inches(8))
    print(f"  Added plot: {plot_path}")
else:
    print(f"  WARNING: Plot not found: {plot_path}")

# Add summary
summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.8))
tf = summary_box.text_frame
tf.text = "Summary: 10/22 achievements unlocked (45%) | Crafter Score: 32.91% | Gap to Human: 50% vs 33%"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 8: Training Metrics (with plot)
print("Creating Slide 8: Training Metrics...")
slide = add_content_slide(prs, "Training Metrics Deep Dive")

# Add plot
plot_path = "logdir/enhanced-dqn/plots/training_metrics.png"
if os.path.exists(plot_path):
    left = Inches(0.3)
    top = Inches(1.3)
    pic = slide.shapes.add_picture(plot_path, left, top, width=Inches(9.4))
    print(f"  Added plot: {plot_path}")
else:
    print(f"  WARNING: Plot not found: {plot_path}")

# Add caption
caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.6))
tf = caption_box.text_frame
tf.text = "Loss converges to ~0.014 | Q-values stabilize at 3-5 | Training shows stable dynamics"
tf.paragraphs[0].font.size = Pt(12)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 9: Emergent Behaviors
print("Creating Slide 9: Emergent Behaviors...")
slide = add_content_slide(prs, "Emergent Behaviors & Observations")

left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.5))
tf = left_box.text_frame
tf.word_wrap = True

add_bullet_text(tf, "✅ Successful Strategies:", level=0, font_size=18, bold=True)
add_bullet_text(tf, "Resource gathering prioritization", level=1, font_size=14)
add_bullet_text(tf, "Food management when low health", level=1, font_size=14)
add_bullet_text(tf, "Day/night adaptation (avoids zombies)", level=1, font_size=14)
add_bullet_text(tf, "Tool progression (wood tools)", level=1, font_size=14)
add_bullet_text(tf, "Infrastructure placement (tables)", level=1, font_size=14)

tf.add_paragraph()
add_bullet_text(tf, "🔍 Interesting Observations:", level=0, font_size=18, bold=True)
add_bullet_text(tf, "Implicit recipe learning (no reward shaping)", level=1, font_size=14)
add_bullet_text(tf, "Risk avoidance (emerges from sparse penalty)", level=1, font_size=14)
add_bullet_text(tf, "Strategy diversity across seeds", level=1, font_size=14)

right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.5))
tf2 = right_box.text_frame
tf2.word_wrap = True

add_bullet_text(tf2, "❌ Failure Modes:", level=0, font_size=18, bold=True)
add_bullet_text(tf2, "Gets stuck in corners (~10%)", level=1, font_size=14)
add_bullet_text(tf2, "Doesn't prioritize rare achievements", level=1, font_size=14)
add_bullet_text(tf2, "Struggles with multi-step planning", level=1, font_size=14)
add_bullet_text(tf2, "Stone/iron tools require 5+ steps", level=1, font_size=14)
add_bullet_text(tf2, "Limited cave exploration", level=1, font_size=14)

# Slide 10: Conclusions
print("Creating Slide 10: Conclusions...")
slide = add_content_slide(prs, "Conclusions & Future Work")

left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.5))
tf = left_box.text_frame
tf.word_wrap = True

add_bullet_text(tf, "✅ Key Achievements:", level=0, font_size=18, bold=True)
add_bullet_text(tf, "Custom DQN from scratch", level=1, font_size=14)
add_bullet_text(tf, "4.90 vs -0.90 reward (5.4× better)", level=1, font_size=14)
add_bullet_text(tf, "10/22 achievements unlocked", level=1, font_size=14)
add_bullet_text(tf, "Double DQN + Dueling = optimal", level=1, font_size=14)
add_bullet_text(tf, "Stable across 3 seeds (±0.16)", level=1, font_size=14)

tf.add_paragraph()
add_bullet_text(tf, "💡 Key Insights:", level=0, font_size=18, bold=True)
add_bullet_text(tf, "Exploration: 50k decay well-calibrated", level=1, font_size=14)
add_bullet_text(tf, "Longer exploration decreased performance", level=1, font_size=14)
add_bullet_text(tf, "Credit assignment = key challenge", level=1, font_size=14)
add_bullet_text(tf, "Simpler can be better", level=1, font_size=14)

right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.5))
tf2 = right_box.text_frame
tf2.word_wrap = True

add_bullet_text(tf2, "🚀 Future Directions:", level=0, font_size=18, bold=True)
add_bullet_text(tf2, "Better Exploration:", level=1, font_size=15, bold=True)
add_bullet_text(tf2, "Intrinsic motivation (RND, NGU)", level=2, font_size=13)
add_bullet_text(tf2, "NOT longer epsilon-greedy", level=2, font_size=13)
tf2.add_paragraph()
add_bullet_text(tf2, "Hierarchical RL:", level=1, font_size=15, bold=True)
add_bullet_text(tf2, "Reusable sub-policies", level=2, font_size=13)
tf2.add_paragraph()
add_bullet_text(tf2, "Model-Based RL:", level=1, font_size=15, bold=True)
add_bullet_text(tf2, "World models (DreamerV2)", level=2, font_size=13)
tf2.add_paragraph()
add_bullet_text(tf2, "Curriculum Learning:", level=1, font_size=15, bold=True)
add_bullet_text(tf2, "Progressive achievement targets", level=2, font_size=13)

# Slide 11: Questions
print("Creating Slide 11: Questions...")
slide = add_content_slide(prs, "Questions & Discussion")

# Thank you box
thank_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1))
tf = thank_box.text_frame
tf.text = "Thank You!"
tf.paragraphs[0].font.size = Pt(48)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Summary box
summary_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2.5))
tf = summary_box.text_frame
tf.word_wrap = True

add_bullet_text(tf, "Summary of Contributions:", level=0, font_size=16, bold=True)
add_bullet_text(tf, "✅ Custom DQN implementation (no frameworks)", level=1, font_size=14)
add_bullet_text(tf, "✅ Systematic ablation study (6 configurations × 3 seeds)", level=1, font_size=14)
add_bullet_text(tf, "✅ Novel exploration experiment (valuable negative result)", level=1, font_size=14)
add_bullet_text(tf, "✅ Comprehensive evaluation (achievements + behaviors)", level=1, font_size=14)

tf.add_paragraph()
p = tf.add_paragraph()
p.text = "Key Takeaway: Systematic investigation beats guesswork"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_ALIGN.CENTER

# Contact box
contact_box = slide.shapes.add_textbox(Inches(2), Inches(6.8), Inches(6), Inches(0.5))
tf = contact_box.text_frame
tf.text = "Contact: [Your Email] | Code: [GitHub URL]"
tf.paragraphs[0].font.size = Pt(12)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save presentation
output_file = "crafter_presentation.pptx"
prs.save(output_file)
print(f"\n✅ Presentation saved as: {output_file}")
print(f"   Total slides: {len(prs.slides)}")
print(f"\nNext steps:")
print(f"  1. Edit '{output_file}' to add your names and contact info")
print(f"  2. Review all slides and adjust formatting as needed")
print(f"  3. Export to PDF: File > Save As > PDF")
